from pptx import Presentation
from pathlib import Path

def extract_text_from_slide(slide):
    text_blocks = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text.strip()
            if text:
                text_blocks.append(text)
    return text_blocks

def pptx_to_qmd(pptx_path="C:/Users/MinYu/pss/PSS - Minyu Chan 27403920.pptx", output_path="C:/Users/MinYu/pss/slides.qmd", title="Converted Slides"):
    # Sanity check: pptx file exists
    if not Path(pptx_path).is_file():
        print(f"❌ File not found: {pptx_path}")
        return

    # Sanity check: output directory exists
    output_dir = Path(output_path).parent
    output_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(pptx_path)
    output_lines = []

    # YAML front matter
    output_lines.append("---")
    output_lines.append(f"title: \"{title}\"")
    output_lines.append("format: revealjs")
    output_lines.append("---\n")

    # Extract each slide
    for idx, slide in enumerate(prs.slides, start=1):
        content = extract_text_from_slide(slide)
        if not content:
            continue  # skip empty slides

        # First block = title, rest = content
        slide_title = content[0] if content else f"Slide {idx}"
        body_lines = content[1:] if len(content) > 1 else []

        output_lines.append(f"# {slide_title}")
        for line in body_lines:
            output_lines.append("")
            output_lines.append(line)

        output_lines.append("\n")  # separate slides

    # Write to file
    Path(output_path).write_text("\n".join(output_lines), encoding="utf-8")
    print(f"✅ Quarto .qmd file written to: {output_path}")

# Example usage:
if __name__ == "__main__":
        pptx_to_qmd(
        pptx_path=r"C:\Users\MinYu\pss\PSS - Minyu Chan 27403920.pptx",
        output_path=r"C:\Users\MinYu\pss\slides.qmd",
        title="Converted Slides"
    )
